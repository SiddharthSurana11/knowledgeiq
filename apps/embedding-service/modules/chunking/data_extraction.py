# modules/data_extraction.py

import os
import logging
from typing import Dict
from dotenv import load_dotenv

import fitz  # PyMuPDF for PDF
import docx
from pptx import Presentation
from PIL import Image
import pytesseract
import io

load_dotenv()
DEBUG_OUTPUT = os.getenv("DEBUG_OUTPUT", "False") == "True"

if os.name == 'nt':
    default_tesseract_path = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    if os.path.exists(default_tesseract_path):
        pytesseract.pytesseract.tesseract_cmd = default_tesseract_path

def extract_text_from_docx(file_path: str) -> str:
    doc = docx.Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text]

    # Extract text from tables
    table_texts = []
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                cell_paragraphs = [p.text for p in cell.paragraphs if p.text]
                table_texts.extend(cell_paragraphs)

    # OCR images in DOCX
    image_texts = []
    rels = doc.part._rels
    for rel in rels:
        rel_obj = rels[rel]
        if "image" in rel_obj.target_ref:
            image_data = rel_obj.target_part.blob
            try:
                image = Image.open(io.BytesIO(image_data))
                text = pytesseract.image_to_string(image)
                if text.strip():
                    image_texts.append(text)
            except Exception as e:
                logging.error(f"OCR failed for image in DOCX: {e}")

    return "\n".join(paragraphs + table_texts + image_texts)

def extract_text_from_pdf(file_path: str) -> Dict[str, any]:
    doc = fitz.open(file_path)
    pages = []
    text_blocks_full = []
    
    for i, page in enumerate(doc):
        text = page.get_text()
        if not text.strip():
            pix = page.get_pixmap()
            img_bytes = pix.tobytes("png")
            image = Image.open(io.BytesIO(img_bytes))
            text = pytesseract.image_to_string(image)
        pages.append({"page": i + 1, "text": text})
        text_blocks_full.append(text)
        
    return {"text": "\n".join(text_blocks_full), "pages": pages}

def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_runs = []
    image_texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
            elif shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    image_bytes = image.blob
                    pil_image = Image.open(io.BytesIO(image_bytes))
                    ocr_text = pytesseract.image_to_string(pil_image)
                    if ocr_text.strip():
                        image_texts.append(ocr_text)
                except Exception as e:
                    logging.error(f"OCR failed for image in PPTX: {e}")
    return "\n".join(text_runs + image_texts)

def extract_text(file_path: str) -> Dict[str, any]:
    """
    Extracts all text (including OCR from embedded images) from PDF, DOCX, PPTX, or images.
    Always returns a dict: { 'text': ..., 'pages': [...], 'filetype': ..., 'error': ... }
    """
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    result: Dict[str, any] = { "text": "", "pages": [], "filetype": ext, "error": "" }
    try:
        if ext == ".pdf":
            pdf_data = extract_text_from_pdf(file_path)
            result["text"] = pdf_data["text"]
            result["pages"] = pdf_data["pages"]
        elif ext in [".docx", ".doc"]:
            result["text"] = extract_text_from_docx(file_path)
        elif ext in [".pptx", ".ppt"]:
            result["text"] = extract_text_from_pptx(file_path)
        elif ext in [".png", ".jpg", ".jpeg"]:
            result["text"] = pytesseract.image_to_string(Image.open(file_path))
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    except Exception as e:
        logging.error(f"Failed to extract text from {file_path}: {e}")
        result["error"] = str(e)
    if DEBUG_OUTPUT and result["text"]:
        base = os.path.basename(file_path)
        with open(f"debug_{base}_extracted.txt", "w", encoding="utf-8") as f:
            f.write(result["text"])
    return result

# Optional: CLI for dev/test (never imported in prod microservice)
if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        f = sys.argv[1]
        result = extract_text(f)
        print(result["text"][:1000])
